"""Build the 2026 seminar deck (16:9) with python-pptx: designed palette, real figures, per-slide speaker script."""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

ROOT = Path(__file__).resolve().parents[2]
FIG = ROOT / "phase1" / "outputs" / "seminar_figures"
PANEL = ROOT / "phase1" / "outputs" / "unet_ds_prior" / "panels"
TEMP = ROOT / "phase1" / "outputs" / "temporal_ds_sequence"
OUTDIR = Path(r"E:/research_projects/Research notes (this folder is a huge mess)")
OUT = OUTDIR / "Spatial Difference Subspaces - seminar 2026 (rebuilt).pptx"

NAVY = RGBColor(0x1D, 0x33, 0x57); TEAL = RGBColor(0x2A, 0x9D, 0x8F)
CORAL = RGBColor(0xE7, 0x6F, 0x51); INK = RGBColor(0x20, 0x25, 0x33)
GRAY = RGBColor(0x5A, 0x64, 0x73); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ICE = RGBColor(0xCA, 0xDC, 0xFC); LIGHT = RGBColor(0xF7, 0xF9, 0xFC)
HEAD, BODY = "Calibri", "Calibri"
EMU = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY if dark else WHITE
    return s


def tb(s, l, t, w, h, lines, size=16, color=INK, bold=False, font=BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        txt, sz, cl, bd = (ln if isinstance(ln, tuple) else (ln, size, color, bold))
        r = p.add_run(); r.text = txt; f = r.font
        f.size = Pt(sz); f.color.rgb = cl; f.bold = bd; f.name = font
    return box


def accent(s, t=0.62):  # teal square motif left of title
    sq = s.shapes.add_shape(1, Inches(0.55), Inches(t + 0.04), Inches(0.16), Inches(0.5))
    sq.fill.solid(); sq.fill.fore_color.rgb = TEAL; sq.line.fill.background()


def title(s, text, sub=None):
    accent(s)
    tb(s, 0.85, 0.5, 11.8, 1.0, text, 30, NAVY, True, HEAD)
    if sub:
        tb(s, 0.85, 1.28, 11.8, 0.5, sub, 16, TEAL, False, BODY)


def footer(s, n, dark=False):
    tb(s, 0.55, 7.04, 9, 0.35, "Abushbeka · Fukui Subspace-Methods Lab, U-Tsukuba", 10, ICE if dark else GRAY)
    tb(s, 12.3, 7.04, 0.6, 0.35, str(n), 10, ICE if dark else GRAY, align=PP_ALIGN.RIGHT)


def pic(s, path, l, t, w=None, h=None, cap=None):
    path = str(path)
    iw, ih = Image.open(path).size; ar = iw / ih
    if w and not h:
        h = w / ar
    elif h and not w:
        w = h * ar
    s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h) if h else None)
    if cap:
        tb(s, l, t + (h or 0) + 0.03, w, 0.35, cap, 11, GRAY, align=PP_ALIGN.CENTER)
    return w, h


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def bullets(items):
    return [("•  " + it, 16, INK, False) for it in items]


# ---------------- SLIDES ----------------
N = 0
def nx():
    global N; N += 1; return N

# 1 title (dark)
s = slide(dark=True)
tb(s, 1.0, 2.5, 11.3, 1.6, "Spatially-faithful Difference Subspaces", 40, WHITE, True, HEAD)
tb(s, 1.0, 3.7, 11.3, 0.9, "for label-free satellite change & disaster-damage detection", 24, ICE, False, HEAD)
tb(s, 1.0, 5.1, 11.3, 0.5, "Abdelrahman I. A. Abushbeka  ·  Kazuhiro Fukui", 18, WHITE)
tb(s, 1.0, 5.6, 11.3, 0.5, "Subspace-Methods Lab, University of Tsukuba", 15, ICE)
sq = s.shapes.add_shape(1, Inches(1.0), Inches(2.25), Inches(2.2), Inches(0.08)); sq.fill.solid(); sq.fill.fore_color.rgb = TEAL; sq.line.fill.background()
notes(s, "Good afternoon. Last year I showed Difference Subspaces and U-Net for change detection on Sentinel-2. This year I went much deeper: I rigorously mapped WHERE subspace geometry genuinely helps satellite change and disaster-damage detection — and I have a real disaster-dataset result. The honest one-line story: geometry is redundant for raw detection, but it adds unique evidence in fusion and it transfers to real disaster damage as a label-free triage tool.")
footer(s, nx(), dark=True)

# 2 outline
s = slide(); title(s, "Outline")
tb(s, 1.0, 2.0, 11, 4.5, bullets([
    "Motivation — fast, label-free disaster damage maps from Sentinel-2",
    "The method — making the Difference Subspace spatially faithful (+ the exact math)",
    "The discipline — naive → complex ladder + matched controls",
    "The diagnostic — where geometry helps vs reduces to a simpler statistic",
    "Results — (A) complementary evidence on OSCD; (B) transfer to real disaster damage (xBD)",
    "What Sensei asked, tried honestly — CCA/S3CCA, kernel, Otsu, geodesic 2nd-order DS",
    "Honest scope, contributions, and future work",
]), 18)
notes(s, "Here's the path: motivation, the method and its math, the rigorous evaluation discipline, the diagnostic that is my backbone, then the two real results, then the methods Sensei specifically asked me to try, and finally honest scope and future work.")
footer(s, nx())

# 3 motivation
s = slide(); title(s, "Motivation: disaster damage mapping", "fast, large-scale, and label-free")
tb(s, 0.85, 2.0, 6.0, 4.3, bullets([
    "After a disaster, responders need damage maps over huge areas — fast.",
    "Ground survey is slow and dangerous.",
    "Sentinel-2: free, global, 13 bands, ~5-day revisit — the practical sensor.",
    "A NEW event has no labels → supervised deep models can't be trained in time.",
    "So the operational tool is UNSUPERVISED change detection.",
]), 17)
pic(s, PANEL / "panel_montpellier.png", 7.0, 2.2, w=5.9, cap="real Sentinel-2 before / after / change (OSCD Montpellier)")
notes(s, "Why this matters: after an earthquake, flood, or conflict, people need to know what is damaged, fast, across a whole region. Ground survey is too slow. Sentinel-2 gives free global imagery every few days, but a brand-new event has no labels — so you cannot train a supervised deep model in time. That forces us to unsupervised change detection. My deep motivation is post-disaster reconstruction.")
footer(s, nx())

# 4 problem
s = slide(); title(s, "Why it's hard — and Sensei's objection")
tb(s, 0.85, 2.0, 11.5, 4.3, bullets([
    "Naive 'spot-the-difference' fails: EVERYTHING changes between two dates (sun, season, registration).",
    "Classical unsupervised maps (CVA, PCA-diff, IR-MAD) are PER-PIXEL → spatially blind.",
    "Our lab's instinct is subspace geometry (Difference Subspace) — but the textbook DS is ALSO per-pixel.",
    "Sensei's objection (last year): \"your algorithm can make a subspace, but it can break the spatial information.\"",
    "→ This year's pivot: make the Difference Subspace spatially faithful, then test honestly where it helps.",
]), 17)
notes(s, "The core difficulty: you can't just subtract two images, because everything changes between two days. The classical unsupervised maps treat each pixel independently, so they're spatially blind. Our lab's instinct is subspace geometry — but the standard Difference Subspace is also per-pixel, which is exactly what Sensei flagged last year: the algorithm makes a subspace but breaks the spatial information. So this year I rebuilt the construction to be spatially faithful, and then — crucially — I tested honestly where it actually helps.")
footer(s, nx())

# 5 framework
s = slide(); title(s, "The proposed framework")
pic(s, FIG / "fig1c_framework.png", 1.35, 1.6, w=10.65)
notes(s, "Here's the pipeline on a real example. Instead of treating each pixel as a sample (which loses position), each whole BAND-IMAGE is one sample — so the subspace lives in spatial coordinates and preserves layout. We build a subspace per date, take their Difference Subspace — the directions where the two dates disagree — and score each pixel by how much its before-after change lives in those directions. This is senpai Jang's 'flatten on the Z-axis' idea.")
footer(s, nx())

# 6 math
s = slide(); title(s, "The math (so I can defend every step)")
pic(s, FIG / "fig9_math_reference.png", 1.6, 1.7, w=10.1)
notes(s, "Precisely, for each date the band-images form a matrix X_t with N pixels by B bands. PCA gives a rank-12 subspace Phi and Psi. Their canonical angles come from the SVD of Phi-transpose-Psi; cosine theta near one means a shared direction. The Difference Subspace D collects the directions that changed. The per-pixel score is the energy of the band-difference projected onto D. Dimension means number of principal components — twelve — not the ambient size. For xBD I also use the projector distance: how much the subspace rotated.")
footer(s, nx())

# 7 discipline (ladder)
s = slide(); title(s, "The discipline: naive → complex + matched controls")
pic(s, FIG / "fig2_ladder.png", 0.7, 1.9, w=8.3)
tb(s, 9.2, 2.1, 3.7, 4.2, bullets([
    "Start from the simplest method (raw difference).",
    "Climb one rung at a time.",
    "Each claim measured vs the simplest thing that could explain it.",
    "Matched 'look-alike' controls isolate what is truly DS-specific.",
]), 15)
notes(s, "Sensei's standing requirement: start from the most naive method and climb complexity one rung at a time. Raw differencing is the floor — the worst. We climb through PCA, smoothed PCA, IR-MAD, then geometry, then the learned U-Net. And we use matched 'look-alike' controls so that any claimed benefit is genuinely from the subspace, not from adding just any extra channel. This discipline is what makes the results trustworthy.")
footer(s, nx())

# 8 diagnostic
s = slide(); title(s, "The diagnostic — the rigorous backbone")
pic(s, FIG / "fig8_diagnostic_matrix.png", 0.7, 1.9, w=11.9)
notes(s, "This is my main contribution as a characterization. A change map is a representation times an operator times a decision. For most change types the fancy geometric operator collapses to a simple one: brightness change is captured by CVA; low-dimensional spectral change makes the DS equal to the spectral angle; change spread across bands is captured by a covariance statistic like IR-MAD; change type is captured by the change direction. Geometry only adds value when the data object is natively curved — like radar covariance matrices — which optical reflectance is not. So this is an honest map of when geometry helps versus when it's redundant — not a claim that it fails.")
footer(s, nx())

# 9 result A OSCD
s = slide(); title(s, "Result A — complementary evidence on OSCD")
pic(s, FIG / "fig3_oscd_fusion_dsspecific.png", 0.7, 1.9, w=8.1)
tb(s, 9.0, 2.1, 4.0, 4.2, bullets([
    "A single DS channel is redundant to a trained CNN.",
    "But in FUSION it adds DS-specific evidence:",
    "0.513 with DS  vs  0.480 without  vs  0.487 look-alike.",
    "Score-level significance p = 0.0098.",
    "Redundant alone — valuable in concert.",
]), 15)
notes(s, "First real result, on the OSCD urban-change benchmark with a U-Net. A single Difference-Subspace channel doesn't help a well-trained CNN — the network learns it. But when we fuse it with the classical maps, it adds evidence that a matched look-alike control cannot reproduce: accuracy rises from 0.48 to 0.51, and at the score level this is significant at p equals 0.0098. So the subspace is redundant alone but valuable in concert.")
footer(s, nx())

# 10 result B xBD headline
s = slide(); title(s, "Result B — it transfers to real disaster damage (xBD-S12)")
pic(s, FIG / "fig4_xbd_headline.png", 0.6, 1.85, w=12.1)
notes(s, "The headline. We froze a protocol and tested on xBD-S12: real Sentinel-2 before-after pairs over disasters, with building-damage masks, on five completely unseen disaster events. The spatial subspace geometry transfers: it beats PCA-diff on all five events and beats the field-standard IR-MAD at localizing buildings. Operationally — the right panel — reviewing just the top five percent of flagged pixels retrieves about a quarter of all damaged pixels, a five-times lift over random. That is a label-free triage tool for disaster response.")
footer(s, nx())

# 11 xBD qualitative
s = slide(); title(s, "Disaster triage in action", "top-5% flagged pixels vs actual damage")
pic(s, FIG / "fig4b_xbd_top5_overlay.png", 0.7, 2.0, w=11.9)
notes(s, "Here it is on one hurricane patch. Left: the after image. Middle: the actual damaged buildings. Right: the top five percent of pixels our method flags for a human to review, in red — they land on the damage clusters, catching 44 percent of the damage here. No labels were used to produce this map. This is the disaster-reconstruction use case that motivates the whole project.")
footer(s, nx())

# 12 honest scope
s = slide(); title(s, "Honest scope — what I do NOT claim")
tb(s, 0.85, 2.0, 11.5, 4.3, bullets([
    "Not state-of-the-art damage assessment, and not pixel-accurate damage segmentation.",
    "Geometry is a candidate-localization prior, not a damage-severity classifier (raw L2 is better at severity).",
    "Only 5 independent test events → consistent directional evidence (p = 0.0625 floor), not p < 0.05.",
    "DS-specificity is strongest at full data; it is budget-dependent.",
    "This honesty is the point: a rigorous, novel first trial — exactly Sensei's bar.",
]), 17)
notes(s, "I want to be precise about what I do not claim. This is not state-of-the-art damage assessment and not pixel-accurate segmentation. The geometry is a candidate-localization prior, not a severity classifier. With only five independent test events, the smallest possible p-value is 0.0625, so I report consistent directional evidence, not statistical significance. Being this honest is the point — Sensei rewards a novel, honest first trial, not an overclaim that collapses under scrutiny.")
footer(s, nx())

# 13 sensei methods (CCA + Otsu)
s = slide(); title(s, "Methods Sensei asked me to try — done")
pic(s, FIG / "fig5_cca_family.png", 0.6, 1.9, w=6.2)
pic(s, FIG / "fig6_otsu_ablation.png", 6.95, 1.9, w=6.0)
tb(s, 0.85, 6.35, 11.6, 0.6, "CCA family (incl. sparse-CCA = the S3CCA flavor) + kernel-DS all tested; trained the U-Net on raw bands with NO Otsu (best), and tested Otsu (far worse).", 13, GRAY)
notes(s, "Sensei asked me to look into specific methods, so I ran them. Left: the CCA family — IR-MAD is iteratively-reweighted CCA, and I added a sparse-CCA which is the S3CCA flavor; the simple CVA still wins here and the subspace is weakest. I also ran the kernel-DS trick — it stays below CVA. Right: the thresholding question — I trained the U-Net on raw bands with no Otsu, which is best, and I tested Otsu thresholding, which roughly halves the score. So every method he named was tried and reported honestly.")
footer(s, nx())

# 14 temporal geodesic
s = slide(); title(s, "Temporal 1st/2nd Difference Subspace + geodesic", "first satellite trial of Sensei's 2024 second-order DS")
pic(s, TEMP / "seq_al_wakrah.png", 1.2, 2.05, w=10.9)
notes(s, "This is the direction Sensei asked about most, and it uses his own 2024 second-order-DS paper. On a real Sentinel-2 time series I compute the first-order DS magnitude — the speed of change — and the second-order DS with its geodesic split into smooth seasonal drift versus abrupt events. It's the first time this is applied to satellite imagery. As Sensei framed it, this is a first and unique trial — a characterization tool, not a detector — and the DS magnitude carries temporal structure beyond a trivial mean shift.")
footer(s, nx())

# 15 interpretability
s = slide(); title(s, "Interpretability, tested honestly")
pic(s, FIG / "fig7_changetype_interpretability.png", 0.7, 1.9, w=7.6)
tb(s, 8.5, 2.1, 4.4, 4.2, bullets([
    "Can the DS directions tell change TYPES apart?",
    "Tested on real 6-class change ground truth (Benton).",
    "Simple change direction (polar-CVA): 0.83.",
    "DS directions: 0.47 — geometry loses.",
    "We checked, and we report it.",
]), 15)
notes(s, "I also tested whether the subspace gives unique interpretable information — specifically, can its directions tell change TYPES apart? On a real six-class change-type ground truth, the simple change direction recovers the types far better than the DS directions: 0.83 versus 0.47. So geometry loses even here. I show this because rigor is credibility — I checked the obvious interpretability hope and report the honest negative.")
footer(s, nx())

# 16 contributions + future
s = slide(); title(s, "Contributions & future work")
tb(s, 0.85, 1.85, 6.0, 4.6, [("Contributions", 19, TEAL, True)] + bullets([
    "A diagnostic: when/why subspace geometry is redundant for optical CD.",
    "DS-specific complementary evidence in fusion (OSCD, p=0.0098).",
    "Transfer to real disaster damage as a label-free candidate-localization prior (xBD).",
]), 15)
tb(s, 7.1, 1.85, 5.7, 4.6, [("Future work", 19, CORAL, True)] + bullets([
    "More independent disaster events (lift p below 0.05).",
    "Neural prior with a fixed projector channel.",
    "Manifold-valued (Grassmann) fusion layers.",
    "13-band GEE temporal sequences for the geodesic analysis.",
]), 15)
notes(s, "To summarize the contributions: first, the diagnostic — a rigorous map of when and why subspace geometry is redundant for optical change detection. Second, DS-specific complementary evidence in fusion on OSCD. Third, transfer to real disaster damage as a label-free candidate-localization prior on xBD. Future work: more disaster events to reach significance, a neural prior using a fixed projector channel, manifold-valued fusion layers, and richer 13-band time series for the geodesic analysis.")
footer(s, nx())

# 17 key references
s = slide(); title(s, "Key references")
refs_l = [
    "Fukui & Maki. Difference Subspace & its generalization. IEEE TPAMI, 2015.",
    "Fukui et al. Second-order Difference Subspace. arXiv:2409.08563, 2024.",
    "Kanai et al. Time-series anomaly detection via DS. arXiv:2303.17802, 2023.",
    "Kobayashi. S3CCA: smoothly structured sparse CCA. ICPR, 2014.",
    "Daudt et al. Fully convolutional Siamese CD (FC-EF). ICIP, 2018.",
    "Nielsen. Regularized iteratively-reweighted MAD (IR-MAD). IEEE TIP, 2007.",
    "Celik. Unsupervised CD via PCA + k-means. IEEE GRSL, 2009.",
    "Bovolo & Bruzzone. CVA in the polar domain. IEEE TGRS, 2007.",
]
refs_r = [
    "Daudt et al. OSCD dataset. IGARSS, 2018.",
    "Gupta et al. xBD building-damage dataset. CVPRW, 2019 (arXiv:1911.09296).",
    "Dietrich/Hafner et al. xBD-S12. arXiv:2511.05461, 2026.",
    "Wenger et al. MultiSenGE Sentinel-2 benchmark, 2022.",
    "Ronneberger et al. U-Net. MICCAI, 2015.",
    "He et al. Deep residual learning (ResNet). CVPR, 2016.",
    "Han et al. Change Guiding Network. arXiv:2404.09179, 2024.",
    "Full list: docs/research/REFERENCES.md",
]
tb(s, 0.85, 1.95, 6.0, 4.8, [(r, 12, INK, False) for r in refs_l], 12, space=8)
tb(s, 7.0, 1.95, 5.9, 4.8, [(r, 12, INK, False) for r in refs_r], 12, space=8)
notes(s, "These are the key references: the lab's subspace lineage (Fukui-Maki TPAMI 2015, the second-order DS and Kanai anomaly papers, S3CCA); the change-detection baselines (FC-Siamese, IR-MAD, PCA-diff, polar CVA); the datasets (OSCD, xBD, xBD-S12, MultiSenGE); and the deep backbones. The full grouped list with links is in REFERENCES.md.")
footer(s, nx())

# 18 conclusion (dark)
s = slide(dark=True)
tb(s, 1.0, 2.4, 11.3, 1.2, "Where subspace geometry genuinely helps", 34, WHITE, True, HEAD)
tb(s, 1.0, 3.7, 11.3, 1.8, [
    ("Redundant for raw detection and change-typing — a simpler statistic wins (the diagnostic).", 18, ICE, False),
    ("Valuable as label-free complementary evidence in fusion, and it transfers to real disaster damage.", 18, ICE, False),
    ("A novel, honest first trial — built to be defended, not overclaimed.", 18, WHITE, True),
], space=10)
tb(s, 1.0, 5.9, 11.3, 0.6, "Thank you — questions welcome.", 20, TEAL, True)
notes(s, "To conclude: subspace geometry is redundant for raw detection and change-typing, where a simpler statistic wins — that's the diagnostic. But it's genuinely valuable as label-free complementary evidence in fusion, and it transfers to real disaster damage. It's a novel, honest first trial, built to be defended rather than overclaimed. Thank you — I'm happy to take questions, and I can derive the difference-subspace construction on the board.")
footer(s, nx(), dark=True)

OUT.parent.mkdir(parents=True, exist_ok=True)
target = OUT
try:
    prs.save(str(target))
except PermissionError:
    target = OUT.with_name(OUT.stem + " (v2).pptx")
    prs.save(str(target))
    print("[note] original was open/locked; saved to v2 instead")
print(f"saved deck: {target}  ({len(prs.slides._sldIdLst)} slides)")
