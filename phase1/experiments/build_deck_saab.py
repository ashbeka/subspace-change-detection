"""Build the 2026 seminar deck (Successive Saab-DS centerpiece, 16:9, ~20 min) with python-pptx."""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

ROOT = Path(__file__).resolve().parents[2]
FIG = ROOT / "phase1" / "outputs" / "seminar_figures"
PANEL = ROOT / "phase1" / "outputs" / "unet_ds_prior" / "panels"
TEMP = ROOT / "phase1" / "outputs" / "temporal_ds_sequence"
CODEX = ROOT / "docs" / "experiment_reports" / "assets" / "multiresolution_subspace_2026-06-23"
OUT = Path(r"E:/research_projects/Research notes (this folder is a huge mess)") / "Spatial Difference Subspaces - seminar 2026 (Saab-DS).pptx"

NAVY = RGBColor(0x16, 0x20, 0x2E); TEAL = RGBColor(0x1D, 0x8F, 0x6E)
CORAL = RGBColor(0xC0, 0x39, 0x2B); INK = RGBColor(0x20, 0x25, 0x33)
GRAY = RGBColor(0x5A, 0x64, 0x73); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ICE = RGBColor(0xC9, 0xE6, 0xDC); LIGHT = RGBColor(0xF7, 0xF9, 0xFB)
HEAD = BODY = "Calibri"

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
N = [0]


def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY if dark else WHITE
    return s


def tb(s, l, t, w, h, lines, size=16, color=INK, bold=False, font=BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        txt, sz, cl, bd = (ln if isinstance(ln, tuple) else (ln, size, color, bold))
        r = p.add_run(); r.text = txt; f = r.font
        f.size = Pt(sz); f.color.rgb = cl; f.bold = bd; f.name = font
    return box


def accent(s, t=0.62):
    sq = s.shapes.add_shape(1, Inches(0.55), Inches(t + 0.04), Inches(0.16), Inches(0.5))
    sq.fill.solid(); sq.fill.fore_color.rgb = TEAL; sq.line.fill.background()


def title(s, text, sub=None):
    accent(s); tb(s, 0.85, 0.5, 11.9, 1.0, text, 29, NAVY, True, HEAD)
    if sub:
        tb(s, 0.85, 1.26, 11.9, 0.5, sub, 16, TEAL, False, BODY)


def footer(s, dark=False):
    N[0] += 1
    tb(s, 0.55, 7.05, 9, 0.35, "Abushbeka · Fukui Subspace-Methods Lab, U-Tsukuba", 10, ICE if dark else GRAY)
    tb(s, 12.3, 7.05, 0.6, 0.35, str(N[0]), 10, ICE if dark else GRAY, align=PP_ALIGN.RIGHT)


def pic(s, path, l, t, w=None, h=None, cap=None):
    path = Path(path)
    if not path.exists():
        print(f"  [skip missing {path.name}]"); return None
    iw, ih = Image.open(path).size; ar = iw / ih
    if w and not h:
        h = w / ar
    elif h and not w:
        w = h * ar
    s.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w), Inches(h))
    if cap:
        tb(s, l, t + h + 0.03, w, 0.35, cap, 11, GRAY, align=PP_ALIGN.CENTER)
    return w, h


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def B(items):
    return [("•  " + it, 16, INK, False) for it in items]


# 1 Title
s = slide(dark=True)
tb(s, 1.0, 2.25, 11.3, 1.6, "Spatially-faithful subspace construction", 36, WHITE, True, HEAD)
tb(s, 1.0, 3.35, 11.3, 1.4, "for unsupervised multispectral change detection — from global pixels to successive Saab features", 21, ICE, False, HEAD)
tb(s, 1.0, 5.2, 11.3, 0.5, "Abdelrahman I. A. Abushbeka  ·  Kazuhiro Fukui", 18, WHITE)
tb(s, 1.0, 5.68, 11.3, 0.5, "Subspace-Methods Lab, University of Tsukuba", 15, ICE)
sq = s.shapes.add_shape(1, Inches(1.0), Inches(2.0), Inches(2.2), Inches(0.08)); sq.fill.solid(); sq.fill.fore_color.rgb = TEAL; sq.line.fill.background()
notes(s, "My project: compare two satellite images, before and after, and automatically find what changed, without any human labels. The key question turns out to be how you turn an image into a subspace. I'll show a journey: the naive way fails exactly as my advisor warned, a spatially-faithful way is better, and a Green-Learning-style successive construction finally beats the classical unsupervised change detectors on held-out cities.")
footer(s, True)

# 2 Outline
s = slide(); title(s, "Outline")
tb(s, 1.0, 1.9, 11, 4.7, B([
    "Motivation — fast, label-free change/damage maps from Sentinel-2",
    "Background — what a subspace is, in one minute",
    "The construction journey — three steps that decide whether geometry helps",
    "The result — successive Saab features + Difference Subspace beat the baselines",
    "Methods Sensei asked me to try (some worked, some did not)",
    "The diagnostic, the temporal extension, and disaster transfer",
    "Honest limits, contributions, and future work",
]), 18)
notes(s, "I'll motivate the problem, define the few terms you need, then walk the construction journey that is the heart of the talk, show the result, show the methods I was asked to try, the temporal extension and disaster relevance, and finish with honest limits and contributions.")
footer(s)

# 3 Motivation
s = slide(); title(s, "Motivation: label-free change & damage mapping")
tb(s, 0.85, 1.95, 6.0, 4.4, B([
    "After a disaster, responders need change/damage maps fast, over huge areas.",
    "Sentinel-2: free, global, 13 bands (incl. infrared), ~5-day revisit.",
    "Each pixel = a vector of 13 numbers (a spectral fingerprint).",
    "A brand-new event has NO labels — can't train a supervised model in time.",
    "So the operational tool is unsupervised (label-free) change detection.",
]), 16)
pic(s, PANEL / "panel_montpellier.png", 7.05, 2.25, w=5.85, cap="real Sentinel-2: before / after / change (OSCD Montpellier)")
notes(s, "Background first. A normal photo has three channels, red green blue. A Sentinel-2 image, which is free and global, has thirteen bands, including infrared that reveals vegetation, water, moisture. So each pixel is a vector of thirteen numbers. The problem: after a disaster, responders need maps fast, and a brand-new event has no labels, so a supervised network can't be trained in time. That forces label-free methods that work straight from the two images.")
footer(s)

# 4 Why hard
s = slide(); title(s, "Why it's hard — and the spatial-information problem")
tb(s, 0.85, 1.95, 11.6, 4.4, B([
    "Naive subtraction fails: everything changes between two dates (sun, season, misalignment).",
    "To use geometry, we turn each image into a 'subspace' — a compact summary via PCA.",
    "But the naive subspace treats pixels as an unordered bag → it loses WHERE things are.",
    "Sensei (last year): \"your algorithm makes a subspace, but it breaks the spatial information.\"",
    "→ This talk fixes that, step by step.",
]), 17)
notes(s, "Why not just subtract? Because everything changes between two days, so subtraction lights up the whole image. Our lab uses subspace geometry: you summarize each image by the main directions its data spread along, found by PCA. But the catch my advisor pointed out: if each sample is one pixel, PCA treats all pixels as an unordered bag and throws away where each pixel is. The subspace knows the colors but not the map. His words were, your algorithm makes a subspace but breaks the spatial information. Fixing that is the spine of this talk.")
footer(s)

# 5 Journey overview
s = slide(); title(s, "The construction journey — the heart of the talk")
pic(s, FIG / "fig_saab_journey.png", 0.7, 1.95, w=11.9)
notes(s, "Here's the whole story in one picture. How you define the sample decides everything. Step one, the naive global-pixel subspace, is spatially blind and scores almost nothing. Step two, the band-image construction, keeps spatial layout and does much better. Step three, adding a label-free Green-Learning successive feature stage before the Difference Subspace, finally beats the classical detectors. I'll unpack each step.")
footer(s)

# 6 Steps 1 & 2
s = slide(); title(s, "Steps 1 & 2: from spatially blind to spatially faithful")
tb(s, 0.85, 1.95, 5.9, 4.4, [("Step 1 — Global pixel DS (naive)", 18, GRAY, True)] + B([
    "sample = one 13-band pixel → subspace in colour space",
    "PCA ignores position → spatially blind",
    "OSCD result: AP 0.06 (far below simple baselines)",
]), 15)
tb(s, 7.0, 1.95, 5.9, 4.4, [("Step 2 — Band-image DS (the fix)", 18, RGBColor(0x2F, 0x5F, 0xA8), True)] + B([
    "sample = one whole band-image → subspace in spatial space",
    "position is preserved (eigen-images)",
    "OSCD result: AP 0.24 — beats geometric nulls, still < smoothed PCA",
]), 15)
tb(s, 0.85, 6.2, 11.6, 0.6, "Progress, not victory — the construction matters, but we still need richer local features.", 14, GRAY)
notes(s, "Step one, the textbook way: each pixel's thirteen numbers is a sample; the Difference Subspace compares before and after. On OSCD this scores Average Precision 0.06, far below even raw subtraction. The spatial-information loss made concrete. Step two, the fix my labmate Jang suggested: flip it, make each whole band-image a sample, so PCA's directions are spatial patterns and position is preserved. This jumps to 0.24 and beats matched geometric controls, but it's still below the best simple method, smoothed PCA. Real progress, but not yet a win. Something is missing: richer local features.")
footer(s)

# 7 Step 3 method
s = slide(); title(s, "Step 3: Successive Saab features + Difference Subspace")
pic(s, FIG / "fig_saab_pipeline.png", 0.45, 1.75, w=12.45)
notes(s, "Step three is the new contribution. Before doing geometry, we build better features with a label-free, training-free stage from Green Learning, called successive Saab. Each hop looks at a pixel's three-by-three neighborhood, so it sees local context, and does a simple PCA that splits the local average, the DC part, from the local patterns, the AC part, keeping about sixteen channels; then it pools and repeats once, so the second hop sees a larger region. The same transform is applied to both dates with no labels. Then we run the Difference Subspace on these features at each scale, and average. And the matched controls at the bottom show the geometry adds beyond the features.")
footer(s)

# 8 Result
s = slide(); title(s, "Result: it beats every classical baseline (frozen, held-out)")
pic(s, FIG / "fig_saab_results.png", 0.6, 1.9, w=7.7)
tb(s, 8.5, 2.0, 4.4, 4.4, B([
    "10 unseen OSCD test cities (frozen).",
    "Saab-DS AP 0.342, AUROC 0.886.",
    "Beats smoothed-PCA 0.314, PCA-diff 0.307, IR-MAD 0.265.",
    "Train-fitted (leakage-free): AP 0.338 → it generalizes.",
    "Seed-stable (3 seeds).",
]), 15)
notes(s, "On ten completely unseen OSCD test cities, frozen, this reaches Average Precision 0.342 and AUROC 0.886, the best of every method we tried, classical or geometric. And one worry I checked directly: was it just adapting to each test pair? No: when I fit the transform only on the training cities and freeze it, it still reaches 0.338 on the unseen test cities. So it's a general representation, not a per-pair trick.")
footer(s)

# 9 Two proofs
s = slide(); title(s, "Why this is solid")
tb(s, 0.85, 1.95, 11.7, 2.4, [("Two findings, backed by matched controls:", 17, TEAL, True)] + B([
    "(1) Successive Saab features are a far better spatial support for DS than global pixels, fixed windows, literal pyramids, or wavelets.",
    "(2) The Difference Subspace adds BEYOND the features: it beats plain L2 and PCA on the SAME features (vs matched PCA +0.037 AP, p=0.002, 10/10 cities).",
]), 15)
pic(s, CODEX / "frozen_test_method_comparison.png", 1.2, 4.3, w=6.2) or pic(s, FIG / "fig_saab_results.png", 3.5, 4.3, w=6.0)
notes(s, "Two findings, and the matched controls make them solid. First, the successive features are a much better input for the subspace than any earlier construction. Second, and this is the key one, the Difference Subspace step adds information beyond the features themselves: on the same Saab features, DS beats plain distance and plain PCA, by 0.037 Average Precision, winning in all ten cities with p equal to 0.002. So it's not just better features; the subspace geometry genuinely contributes.")
footer(s)

# 10 Methods asked
s = slide(); title(s, "Methods Sensei asked me to try (honest outcomes)")
pic(s, FIG / "fig5_cca_family.png", 0.55, 2.0, w=6.2)
pic(s, FIG / "fig6_otsu_ablation.png", 6.95, 2.0, w=6.0)
tb(s, 0.85, 6.35, 11.7, 0.7, "Green Learning/Saab → WORKED (the result). Literal pyramids (2×2/4×4) → 0.21, didn't help. Wavelets (Haar/db2) → 0.21, detail = nuisance. Kernel-DS → below CVA. CCA/S3CCA → tested.", 12.5, GRAY)
notes(s, "My advisor suggested several techniques, so I ran all of them, and I'm honest about the outcome. The Green-Learning successive Saab idea is the one that worked. But the literal multi-scale pyramid, two-by-two and four-by-four, scored 0.21, below plain PCA differencing. And wavelets, both Haar and Daubechies, also 0.21; their high-frequency detail mostly captured misregistration noise, not real change. I also ran the kernel version, still below the simplest baseline, and the CCA family he named. So I did everything he asked, and I can tell him exactly which ideas paid off and which didn't, with numbers.")
footer(s)

# 11 Diagnostic
s = slide(); title(s, "The diagnostic — the honest backbone")
pic(s, FIG / "fig8_diagnostic_matrix.png", 0.7, 1.95, w=11.9)
notes(s, "Underneath this is a general lesson I call the diagnostic: a change detector is a representation, times a comparison, times a decision, and for many change types the fancy geometry collapses into a simpler statistic. That's why naive subspaces lose. The value of this project is showing precisely when the construction, the representation, is what makes geometry finally useful. The journey from 0.06 to 0.34 is that lesson made concrete.")
footer(s)

# 12 Temporal
s = slide(); title(s, "Temporal extension: 1st/2nd DS + geodesic", "first satellite trial of Sensei's 2024 second-order DS")
pic(s, TEMP / "seq_al_wakrah.png", 1.2, 2.05, w=10.9)
notes(s, "My advisor's most repeated request was to extend across time, using his own 2024 second-order Difference-Subspace paper. The same per-date subspaces become points on a curved surface; how fast you move is the first-order magnitude, the speed of change, and how the path bends is the second-order, the acceleration, which separates smooth seasonal drift from sudden events. First time it's applied to satellites. As he framed it, a first and unique trial, a characterization tool, not a detector.")
footer(s)

# 13 Disaster
s = slide(); title(s, "Disaster relevance: candidate localization (xBD)", "top-5% flagged pixels vs actual damage — no labels used")
pic(s, FIG / "fig4b_xbd_top5_overlay.png", 0.7, 2.05, w=11.9)
notes(s, "Does any of this help on real disasters? On xBD-S12, actual hurricanes with human-marked building damage, the spatial-subspace map beats classical maps at locating damaged buildings: if a human reviews only the most-suspicious five percent of pixels, they already find about a quarter of the damage, five times better than random, with no labels. Triage means prioritizing what a person checks first. The full successive method does not automatically transfer here yet, so I present the simpler band-image projector for disaster, and Saab-DS disaster transfer as future work.")
footer(s)

# 14 Honest scope
s = slide(); title(s, "Honest scope — what I do NOT claim")
tb(s, 0.85, 1.95, 11.6, 4.4, B([
    "Not state-of-the-art; does not beat modern supervised deep change detectors (it's label-free).",
    "It generalizes — a transform fit only on training cities reaches AP 0.338 on unseen cities (≈ 0.342).",
    "vs the single strongest baseline (smoothed PCA): a clear trend (8/10 cities) but p=0.084 with only 10 cities; significant vs raw/PCA.",
    "Disaster transfer (xBD) is mixed — Saab-DS doesn't beat simple L2 there; the band-image projector is the disaster tool.",
    "A scoped adaptation of PixelHop/Saab, not the full Green-Learning system.",
]), 16)
notes(s, "Honest limits. Not state-of-the-art, doesn't beat modern supervised networks; it's label-free. It does generalize: a transform fit only on training cities reaches 0.338 on unseen cities. With only ten test cities, the win over the single strongest baseline, smoothed PCA, is a clear trend but not yet formally significant, though it is significant against raw differencing and PCA. On real disaster data this exact method doesn't beat simple methods, so I present the projector for disaster. And it's a scoped adaptation, not the full PixelHop. Being precise here is what makes the positive credible.")
footer(s)

# 15 Contributions
s = slide(); title(s, "Contributions & future work")
tb(s, 0.85, 1.85, 6.1, 4.7, [("Contributions", 19, TEAL, True)] + B([
    "The construction journey: how the satellite sample definition decides whether DS helps.",
    "Successive Saab + DS beats classical unsupervised detectors on held-out OSCD; DS adds beyond the features.",
    "Honest negatives for literal pyramids, wavelets, and kernel-DS.",
    "The diagnostic + the first satellite trial of 2nd-order DS / geodesic.",
]), 14)
tb(s, 7.1, 1.85, 5.7, 4.7, [("Future work", 19, CORAL, True)] + B([
    "More test cities / a 2nd labeled benchmark → formal significance.",
    "Transfer Saab-DS to disaster damage (xBD-S12).",
    "Use the map as a label-efficient prior for a network.",
    "Multi-date temporal Saab-DS + geodesic analysis.",
]), 14)
notes(s, "Contributions: first, the construction journey, showing how the choice of sample decides whether geometry helps. Second, a concrete method, successive Saab plus Difference Subspace, that beats the classical unsupervised detectors on held-out cities, with the geometry adding beyond the features. Third, clean negatives for pyramids, wavelets, and kernel. Fourth, the diagnostic and the first satellite trial of the second-order DS. Future work: more cities for significance, disaster transfer, a label-efficient neural prior, and multi-date temporal analysis.")
footer(s)

# 16 References
s = slide(); title(s, "Key references")
refs_l = [
    "Fukui & Maki. Difference Subspace & its generalization. IEEE TPAMI, 2015.",
    "Fukui et al. Second-order Difference Subspace. arXiv:2409.08563, 2024.",
    "Kuo et al. PixelHop: Successive Subspace Learning. arXiv:1909.08190, 2019.",
    "Kuo & Madni. Green Learning: introduction & outlook. JVCIR, 2022.",
    "Daudt et al. Fully convolutional Siamese CD (FC-EF). ICIP, 2018.",
]
refs_r = [
    "Nielsen. IR-MAD change detection. IEEE TIP, 2007.",
    "Celik. Unsupervised CD via PCA + k-means. IEEE GRSL, 2009.",
    "Bovolo & Bruzzone. CVA in the polar domain. IEEE TGRS, 2007.",
    "Daudt et al. OSCD dataset. IGARSS, 2018.",
    "Mallat. Wavelet multiresolution. IEEE TPAMI, 1989.   ·   Full list: REFERENCES.md",
]
tb(s, 0.85, 2.0, 6.1, 4.6, [(r, 12, INK, False) for r in refs_l], 12, space=9)
tb(s, 7.0, 2.0, 5.9, 4.6, [(r, 12, INK, False) for r in refs_r], 12, space=9)
notes(s, "Key references: the lab's subspace papers, the Green-Learning and PixelHop papers behind the Saab features, the classical change-detection baselines, the OSCD dataset, and wavelets. The full grouped list is in my references file.")
footer(s)

# 17 Conclusion
s = slide(dark=True)
tb(s, 1.0, 2.3, 11.3, 1.2, "How you build the subspace decides everything", 33, WHITE, True, HEAD)
tb(s, 1.0, 3.6, 11.3, 2.0, [
    ("A Green-Learning 'successive Saab' construction finally makes Difference-Subspace geometry beat the classical unsupervised detectors on held-out cities — and the geometry adds beyond the features.", 18, ICE, False),
    ("A careful, honest study that does exactly what Sensei asked — including the parts that didn't work.", 18, WHITE, True),
], space=12)
tb(s, 1.0, 5.95, 11.3, 0.6, "Thank you — questions welcome.", 20, TEAL, True)
notes(s, "To conclude: how you turn a satellite image into a subspace decides whether the geometry is useless or useful. A label-free, Green-Learning-style successive construction finally makes the Difference Subspace beat the classical unsupervised detectors on unseen cities, and the geometry adds beyond the features. A careful, honest study that does exactly what my advisor asked, including the parts that didn't work. Thank you. I'm happy to take questions and to derive the construction on the board.")
footer(s, True)

OUT.parent.mkdir(parents=True, exist_ok=True)
try:
    prs.save(str(OUT))
except PermissionError:
    OUT = OUT.with_name(OUT.stem + " v2.pptx"); prs.save(str(OUT))
print(f"saved deck: {OUT}  ({len(prs.slides._sldIdLst)} slides)")
